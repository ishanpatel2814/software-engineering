"""
PowerPoint processor for extracting content from PPT/PPTX files.
Renders real slide images via LibreOffice->PDF and PyMuPDF (fitz).
"""
from __future__ import annotations

import os
import logging
import tempfile
import subprocess
from typing import Dict, List, Any, Tuple

from pptx import Presentation
import fitz  # PyMuPDF

from utils.file_utils import create_temp_directory

logger = logging.getLogger("edu_video_pipeline")


class PPTProcessor:
    """Processor for extracting content from PowerPoint files."""

    def __init__(self, config: Dict[str, Any]):
        self.config = config
        self.temp_dir = create_temp_directory(config.get("TEMP_DIR"))

    # ------- helpers ---------------------------------------------------------

    def _target_size(self) -> Tuple[int, int]:
        res = self.config.get("OUTPUT_RESOLUTION", (1920, 1080))
        try:
            w, h = res
            return int(w), int(h)
        except Exception:
            return (1920, 1080)

    def _ensure_pptx(self, ppt_path: str) -> str:
        """Return a .pptx path. If input is .ppt, convert it using LibreOffice."""
        if ppt_path.lower().endswith(".pptx"):
            return ppt_path

        logger.info(f"Converting legacy .ppt to .pptx: {ppt_path}")
        with tempfile.TemporaryDirectory(dir=self.temp_dir) as tmpd:
            cmd = [
                "libreoffice", "--headless",
                "--convert-to", "pptx",
                "--outdir", tmpd,
                ppt_path,
            ]
            subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            pptx_path = None
            for f in os.listdir(tmpd):
                if f.lower().endswith(".pptx"):
                    pptx_path = os.path.join(tmpd, f)
                    break
            if not pptx_path:
                raise RuntimeError("Failed to convert .ppt to .pptx with LibreOffice.")

            # move into our temp dir for persistence after tmpd is deleted
            dst = os.path.join(self.temp_dir, os.path.basename(pptx_path))
            os.replace(pptx_path, dst)
            return dst

    def _convert_to_pdf(self, in_path: str) -> str:
        """Convert PPT/PPTX to PDF using LibreOffice."""
        outdir = tempfile.mkdtemp(prefix="pptpdf_", dir=self.temp_dir)
        cmd = [
            "libreoffice", "--headless",
            "--convert-to", "pdf",
            "--outdir", outdir,
            in_path,
        ]
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)

        for f in os.listdir(outdir):
            if f.lower().endswith(".pdf"):
                return os.path.join(outdir, f)
        raise RuntimeError("LibreOffice did not produce a PDF.")

    # ------- public API ------------------------------------------------------

    def extract_text(self, ppt_path: str) -> List[Dict[str, Any]]:
        """
        Extract text content per slide using python-pptx.
        (Auto-converts .ppt → .pptx first.)
        """
        text_content: List[Dict[str, Any]] = []
        try:
            pptx_path = self._ensure_pptx(ppt_path)
            presentation = Presentation(pptx_path)

            for slide_num, slide in enumerate(presentation.slides, start=1):
                slide_text = []
                shapes_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        t = shape.text.strip()
                        if t:
                            slide_text.append(t)
                            left = getattr(shape, "left", 0)
                            top = getattr(shape, "top", 0)
                            width = getattr(shape, "width", 0)
                            height = getattr(shape, "height", 0)
                            shapes_text.append({
                                "text": t,
                                "bbox": (left, top, left + width, top + height),
                            })
                text_content.append({
                    "slide_num": slide_num,
                    "text": "\n".join(slide_text).strip(),
                    "shapes": shapes_text,
                })

            logger.info(f"Extracted text from {len(presentation.slides)} slides in PPT: {ppt_path}")
            return text_content

        except Exception as e:
            logger.error(f"Error extracting text from PPT {ppt_path}: {e}")
            raise

    def extract_slides(self, ppt_path: str) -> List[Dict[str, Any]]:
        """
        Render each slide to a sharp PNG by converting to PDF then rasterizing with PyMuPDF.
        """
        slides: List[Dict[str, Any]] = []
        try:
            pdf_path = self._convert_to_pdf(ppt_path)

            base = os.path.splitext(os.path.basename(ppt_path))[0]
            slide_dir = os.path.join(self.temp_dir, f"slides_{base}")
            os.makedirs(slide_dir, exist_ok=True)

            doc = fitz.open(pdf_path)
            target_w, target_h = self._target_size()

            for i, page in enumerate(doc, start=1):
                pw, ph = page.rect.width, page.rect.height  # points
                # zoom so rendered pixels ~= target resolution
                zoom = max(target_w / pw, target_h / ph)
                mat = fitz.Matrix(zoom, zoom)

                pix = page.get_pixmap(matrix=mat, alpha=False)  # RGB
                out_path = os.path.join(slide_dir, f"slide_{i}.png")
                pix.save(out_path)

                slides.append({
                    "slide_num": i,
                    "path": out_path,
                    "width": pix.width,
                    "height": pix.height,
                })

            doc.close()
            logger.info(f"Extracted {len(slides)} slide images from PPT: {ppt_path}")
            return slides

        except FileNotFoundError:
            logger.error(
                "LibreOffice not found. Install it with: sudo apt-get update && sudo apt-get install -y libreoffice"
            )
            raise
        except subprocess.CalledProcessError as e:
            logger.error(f"LibreOffice conversion failed: {e}")
            raise
        except Exception as e:
            logger.error(f"Error extracting slides from PPT {ppt_path}: {e}")
            raise

    def extract_notes(self, ppt_path: str) -> List[Dict[str, Any]]:
        """
        Extract speaker notes per slide.
        """
        notes: List[Dict[str, Any]] = []
        try:
            pptx_path = self._ensure_pptx(ppt_path)
            presentation = Presentation(pptx_path)

            for slide_num, slide in enumerate(presentation.slides, start=1):
                notes_text = ""
                if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text or ""
                notes.append({"slide_num": slide_num, "notes": notes_text.strip()})

            logger.info(f"Extracted notes from {len(presentation.slides)} slides in PPT: {ppt_path}")
            return notes

        except Exception as e:
            logger.error(f"Error extracting notes from PPT {ppt_path}: {e}")
            raise

    def get_metadata(self, ppt_path: str) -> Dict[str, Any]:
        """Return basic metadata about the deck."""
        try:
            pptx_path = self._ensure_pptx(ppt_path)
            presentation = Presentation(pptx_path)
            core = presentation.core_properties
            return {
                "title": core.title or "",
                "author": core.author or "",
                "subject": core.subject or "",
                "keywords": core.keywords or "",
                "created": core.created.isoformat() if core.created else "",
                "modified": core.modified.isoformat() if core.modified else "",
                "slide_count": len(presentation.slides),
                "file_size": os.path.getsize(ppt_path),
            }
        except Exception as e:
            logger.error(f"Error getting metadata from PPT {ppt_path}: {e}")
            raise
